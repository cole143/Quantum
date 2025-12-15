"""
╔═══════════════════════════════════════════════════════════════════════╗
║                                                                       ║
║                     QUANTUM - PYTHON SLIDE GENERATION CODE            ║
║                     Luxury Financial Software Tool                    ║
║                     AI-Powered Presentation Generator                 ║
║                                                                       ║
╚═══════════════════════════════════════════════════════════════════════╝

COMPLETE PYTHON CODE FOR QUANTUM PRESENTATION GENERATION

This file contains all the Python code needed to generate professional,
luxury financial presentations with AI-powered content and brand learning.

CONTENTS:
=========
1. QUANTUM CORE SYSTEM (Lines 30-669)
   - QuantumStyleAnalyzer: Learns formatting from reference decks
   - QuantumPresentation: Generates professional slides
   - QuantumAI: AI-powered content generation
   - Main generation function

2. USAGE EXAMPLES (Lines 671-926)
   - Example 1: Quarterly Financial Report
   - Example 2: Investment Pitch Deck
   - Example 3: Annual Performance Review
   - Example 4: Custom Luxury Styling
   - Example 5: Reference Deck Learning

3. QUICK START SCRIPT (Lines 928-981)
   - Simplest possible usage
   - Immediate results
   - Sample data included

INSTALLATION:
=============
pip install python-pptx pillow

BASIC USAGE:
============
from quantum import create_quantum_presentation

data = {
    'metrics': {'Revenue': '$125M', 'Growth': '+15%'},
    'analysis': {'trend': 'positive'}
}

create_quantum_presentation(
    "Q4 Financial Review",
    data,
    output_filename="report.pptx"
)

WITH BRAND LEARNING:
===================
create_quantum_presentation(
    "Branded Report",
    data,
    reference_deck_path="your_template.pptx",
    output_filename="branded_report.pptx"
)

"""

# ═══════════════════════════════════════════════════════════════════════════
# SECTION 1: QUANTUM CORE SYSTEM
# ═══════════════════════════════════════════════════════════════════════════

"""
Quantum - Luxury Financial Software Tool
Professional AI-Powered Presentation Generator

A sophisticated system for creating premium financial slide decks
with learned formatting standards and AI-driven content generation.
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass, asdict
import re


@dataclass
class SlideStyle:
    """Captures the luxury formatting standards from reference decks"""
    title_font_size: int = 44
    title_font_name: str = "Calibri"
    title_color: Tuple[int, int, int] = (0, 0, 0)
    body_font_size: int = 18
    body_font_name: str = "Calibri"
    body_color: Tuple[int, int, int] = (64, 64, 64)
    background_color: Optional[Tuple[int, int, int]] = None
    accent_color: Tuple[int, int, int] = (0, 102, 204)
    layout_preferences: Dict[str, int] = None
    
    def to_dict(self) -> dict:
        return asdict(self)
    
    @classmethod
    def from_dict(cls, data: dict) -> 'SlideStyle':
        if 'layout_preferences' in data and data['layout_preferences'] is None:
            data['layout_preferences'] = {}
        return cls(**data)


class QuantumStyleAnalyzer:
    """Analyzes reference presentations to learn luxury formatting standards"""
    
    def __init__(self):
        self.style = SlideStyle()
    
    def analyze_reference_deck(self, reference_path: str) -> SlideStyle:
        """
        Analyzes a reference PowerPoint to extract premium styling standards
        
        Args:
            reference_path: Path to the reference .pptx file
            
        Returns:
            SlideStyle object containing learned formatting standards
        """
        try:
            prs = Presentation(reference_path)
            
            # Analyze title slides
            title_fonts = []
            title_sizes = []
            title_colors = []
            
            # Analyze body text
            body_fonts = []
            body_sizes = []
            body_colors = []
            
            # Track layout usage
            layout_count = {}
            
            for slide in prs.slides:
                layout_name = slide.slide_layout.name
                layout_count[layout_name] = layout_count.get(layout_name, 0) + 1
                
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size:
                                    size = run.font.size.pt
                                    font_name = run.font.name or "Calibri"
                                    
                                    # Detect if this is likely a title (larger font)
                                    if size >= 32:
                                        title_sizes.append(size)
                                        title_fonts.append(font_name)
                                        if run.font.color and run.font.color.rgb:
                                            title_colors.append(run.font.color.rgb)
                                    elif size >= 12:
                                        body_sizes.append(size)
                                        body_fonts.append(font_name)
                                        if run.font.color and run.font.color.rgb:
                                            body_colors.append(run.font.color.rgb)
            
            # Calculate most common values
            if title_sizes:
                self.style.title_font_size = int(self._most_common(title_sizes))
            if title_fonts:
                self.style.title_font_name = self._most_common(title_fonts)
            if title_colors:
                common_color = self._most_common(title_colors)
                self.style.title_color = (common_color[0], common_color[1], common_color[2])
            
            if body_sizes:
                self.style.body_font_size = int(self._most_common(body_sizes))
            if body_fonts:
                self.style.body_font_name = self._most_common(body_fonts)
            if body_colors:
                common_color = self._most_common(body_colors)
                self.style.body_color = (common_color[0], common_color[1], common_color[2])
            
            self.style.layout_preferences = layout_count
            
            print(f"✓ Quantum analyzed reference deck: {os.path.basename(reference_path)}")
            print(f"  → Title Style: {self.style.title_font_name}, {self.style.title_font_size}pt")
            print(f"  → Body Style: {self.style.body_font_name}, {self.style.body_font_size}pt")
            print(f"  → Layouts found: {len(layout_count)}")
            
            return self.style
            
        except Exception as e:
            print(f"✗ Error analyzing reference deck: {str(e)}")
            print("  → Using default luxury styling standards")
            return self.style
    
    def _most_common(self, items):
        """Helper to find most common item in a list"""
        if not items:
            return None
        return max(set(items), key=items.count)
    
    def save_style(self, path: str):
        """Save learned style to JSON for reuse"""
        with open(path, 'w') as f:
            json.dump(self.style.to_dict(), f, indent=2)
        print(f"✓ Style profile saved to {path}")
    
    def load_style(self, path: str):
        """Load previously learned style from JSON"""
        try:
            with open(path, 'r') as f:
                data = json.load(f)
                self.style = SlideStyle.from_dict(data)
            print(f"✓ Style profile loaded from {path}")
        except Exception as e:
            print(f"✗ Could not load style: {str(e)}")


class QuantumPresentation:
    """
    Quantum Luxury Financial Presentation Generator
    
    Creates professional, AI-driven financial slide decks with learned
    formatting standards and premium visual design.
    """
    
    def __init__(self, style: Optional[SlideStyle] = None):
        """
        Initialize Quantum with optional custom styling
        
        Args:
            style: SlideStyle object with luxury formatting standards
        """
        self.prs = Presentation()
        self.style = style or SlideStyle()
        self.current_slide = None
        
        # Set presentation dimensions (16:9 widescreen for luxury feel)
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self, title: str, subtitle: str = ""):
        """
        Add a luxury-styled title slide
        
        Args:
            title: Main presentation title
            subtitle: Optional subtitle or tagline
        """
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Style the title
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
            self._apply_title_style(title_shape)
        
        # Style the subtitle
        if len(slide.placeholders) > 1 and subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            self._apply_body_style(subtitle_shape)
        
        self.current_slide = slide
        return slide
    
    def add_content_slide(self, title: str, content: List[str], layout_type: str = "content"):
        """
        Add a content slide with bullet points
        
        Args:
            title: Slide title
            content: List of bullet points or content items
            layout_type: Type of layout ('content', 'two_column', 'blank')
        """
        # Select appropriate layout
        layout_map = {
            "content": 1,
            "two_column": 3,
            "blank": 6
        }
        layout_idx = layout_map.get(layout_type, 1)
        
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add and style title
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._apply_title_style(slide.shapes.title)
        
        # Add content
        if len(slide.placeholders) > 1:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            for idx, item in enumerate(content):
                p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
                p.text = item
                p.level = 0
                self._apply_paragraph_style(p)
        
        self.current_slide = slide
        return slide
    
    def add_financial_data_slide(self, title: str, data: Dict[str, Any], 
                                 chart_type: str = "table"):
        """
        Add a slide with financial data visualization
        
        Args:
            title: Slide title
            data: Dictionary containing financial data
            chart_type: Type of visualization ('table', 'metrics', 'comparison')
        """
        slide_layout = self.prs.slide_layouts[5]  # Blank layout for custom content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
        )
        title_shape.text = title
        self._apply_title_style(title_shape)
        
        # Add financial data based on type
        if chart_type == "metrics":
            self._add_financial_metrics(slide, data)
        elif chart_type == "table":
            self._add_financial_table(slide, data)
        elif chart_type == "comparison":
            self._add_comparison_view(slide, data)
        
        self.current_slide = slide
        return slide
    
    def _add_financial_metrics(self, slide, data: Dict[str, Any]):
        """Add key financial metrics in a luxury grid layout"""
        metrics = data.get('metrics', {})
        
        # Calculate grid positions for up to 6 metrics
        start_top = Inches(2)
        start_left = Inches(1)
        box_width = Inches(3.5)
        box_height = Inches(1.5)
        spacing = Inches(0.3)
        
        items = list(metrics.items())[:6]
        
        for idx, (metric_name, metric_value) in enumerate(items):
            row = idx // 3
            col = idx % 3
            
            left = start_left + col * (box_width + spacing)
            top = start_top + row * (box_height + spacing)
            
            # Create metric box
            box = slide.shapes.add_textbox(left, top, box_width, box_height)
            text_frame = box.text_frame
            
            # Add metric name
            p1 = text_frame.paragraphs[0]
            p1.text = metric_name
            p1.font.size = Pt(self.style.body_font_size - 2)
            p1.font.name = self.style.body_font_name
            p1.font.color.rgb = RGBColor(*self.style.body_color)
            
            # Add metric value
            p2 = text_frame.add_paragraph()
            p2.text = str(metric_value)
            p2.font.size = Pt(self.style.title_font_size - 8)
            p2.font.name = self.style.title_font_name
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(*self.style.accent_color)
    
    def _add_financial_table(self, slide, data: Dict[str, Any]):
        """Add a formatted financial data table"""
        table_data = data.get('table', [])
        
        if not table_data:
            return
        
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        
        if rows == 0 or cols == 0:
            return
        
        # Add table
        left = Inches(1)
        top = Inches(2)
        width = Inches(11.333)
        height = Inches(4.5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Populate and style table
        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_value)
                
                # Style header row
                if row_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*self.style.accent_color)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(self.style.body_font_size)
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                else:
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(self.style.body_font_size - 2)
                        paragraph.font.name = self.style.body_font_name
    
    def _add_comparison_view(self, slide, data: Dict[str, Any]):
        """Add a comparison view for financial data"""
        comparisons = data.get('comparisons', {})
        
        left_col = Inches(1)
        right_col = Inches(7)
        top = Inches(2.5)
        width = Inches(5.5)
        
        for idx, (label, values) in enumerate(comparisons.items()):
            # Left column
            left_box = slide.shapes.add_textbox(
                left_col, top + idx * Inches(1.2), width, Inches(1)
            )
            left_frame = left_box.text_frame
            left_p = left_frame.paragraphs[0]
            left_p.text = f"{label}: {values.get('current', 'N/A')}"
            left_p.font.size = Pt(self.style.body_font_size)
            left_p.font.name = self.style.body_font_name
            
            # Right column (comparison)
            right_box = slide.shapes.add_textbox(
                right_col, top + idx * Inches(1.2), width, Inches(1)
            )
            right_frame = right_box.text_frame
            right_p = right_frame.paragraphs[0]
            change = values.get('change', 'N/A')
            right_p.text = f"Change: {change}"
            right_p.font.size = Pt(self.style.body_font_size)
            right_p.font.name = self.style.body_font_name
            
            # Color code positive/negative
            if isinstance(change, str) and '%' in change:
                if '-' in change:
                    right_p.font.color.rgb = RGBColor(204, 0, 0)  # Red
                else:
                    right_p.font.color.rgb = RGBColor(0, 153, 0)  # Green
    
    def _apply_title_style(self, shape):
        """Apply luxury title formatting"""
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(self.style.title_font_size)
            paragraph.font.name = self.style.title_font_name
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(*self.style.title_color)
            paragraph.alignment = PP_ALIGN.LEFT
    
    def _apply_body_style(self, shape):
        """Apply luxury body text formatting"""
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            self._apply_paragraph_style(paragraph)
    
    def _apply_paragraph_style(self, paragraph):
        """Apply luxury paragraph formatting"""
        paragraph.font.size = Pt(self.style.body_font_size)
        paragraph.font.name = self.style.body_font_name
        paragraph.font.color.rgb = RGBColor(*self.style.body_color)
        paragraph.alignment = PP_ALIGN.LEFT
    
    def save(self, filename: str):
        """
        Save the luxury presentation
        
        Args:
            filename: Output filename for the .pptx file
        """
        self.prs.save(filename)
        print(f"✓ Quantum presentation saved: {filename}")
        return filename


class QuantumAI:
    """
    AI-Powered Content Generator for Quantum
    
    Generates intelligent, context-aware financial content for presentations
    """
    
    @staticmethod
    def generate_executive_summary(financial_data: Dict[str, Any]) -> List[str]:
        """
        Generate executive summary bullet points from financial data
        
        Args:
            financial_data: Dictionary containing financial metrics and analysis
            
        Returns:
            List of executive summary points
        """
        summary = []
        
        metrics = financial_data.get('metrics', {})
        
        # Revenue analysis
        if 'revenue' in metrics:
            revenue = metrics['revenue']
            summary.append(f"Total revenue of {revenue} demonstrating strong market performance")
        
        # Profitability
        if 'profit_margin' in metrics:
            margin = metrics['profit_margin']
            summary.append(f"Profit margin of {margin} reflecting operational efficiency")
        
        # Growth metrics
        if 'growth_rate' in metrics:
            growth = metrics['growth_rate']
            summary.append(f"Year-over-year growth rate of {growth} indicating positive momentum")
        
        # Market position
        if 'market_share' in metrics:
            share = metrics['market_share']
            summary.append(f"Market share of {share} establishing competitive positioning")
        
        # Future outlook
        summary.append("Strategic initiatives aligned for continued value creation")
        
        return summary
    
    @staticmethod
    def generate_recommendations(analysis: Dict[str, Any]) -> List[str]:
        """Generate strategic recommendations based on financial analysis"""
        recommendations = []
        
        if analysis.get('trend') == 'positive':
            recommendations.append("Continue current growth trajectory with strategic investments")
            recommendations.append("Expand market presence in high-performing segments")
        elif analysis.get('trend') == 'negative':
            recommendations.append("Implement cost optimization strategies")
            recommendations.append("Reassess market positioning and competitive advantages")
        else:
            recommendations.append("Maintain operational efficiency while exploring new opportunities")
        
        recommendations.append("Enhance stakeholder communication and transparency")
        recommendations.append("Monitor key performance indicators for early trend detection")
        
        return recommendations
    
    @staticmethod
    def format_currency(value: float, currency: str = "USD") -> str:
        """Format numerical values as currency"""
        if currency == "USD":
            if value >= 1_000_000_000:
                return f"${value/1_000_000_000:.2f}B"
            elif value >= 1_000_000:
                return f"${value/1_000_000:.2f}M"
            elif value >= 1_000:
                return f"${value/1_000:.2f}K"
            else:
                return f"${value:.2f}"
        return f"{value:,.2f}"
    
    @staticmethod
    def calculate_growth(current: float, previous: float) -> str:
        """Calculate and format growth percentage"""
        if previous == 0:
            return "N/A"
        growth = ((current - previous) / previous) * 100
        sign = "+" if growth >= 0 else ""
        return f"{sign}{growth:.1f}%"


def create_quantum_presentation(
    title: str,
    financial_data: Dict[str, Any],
    reference_deck_path: Optional[str] = None,
    output_filename: str = "quantum_presentation.pptx"
) -> str:
    """
    Main function to create a complete Quantum luxury financial presentation
    
    Args:
        title: Presentation title
        financial_data: Dictionary containing all financial data and metrics
        reference_deck_path: Optional path to reference deck for style learning
        output_filename: Output filename for the presentation
        
    Returns:
        Path to the generated presentation file
    """
    print("=" * 60)
    print("QUANTUM - Luxury Financial Software Tool")
    print("=" * 60)
    
    # Initialize style analyzer
    analyzer = QuantumStyleAnalyzer()
    
    # Learn from reference deck if provided
    if reference_deck_path and os.path.exists(reference_deck_path):
        style = analyzer.analyze_reference_deck(reference_deck_path)
    else:
        style = SlideStyle()
        print("→ Using default luxury styling standards")
    
    # Initialize AI and presentation
    ai = QuantumAI()
    prs = QuantumPresentation(style)
    
    # Create title slide
    print("\n→ Generating title slide...")
    prs.add_title_slide(
        title,
        "Premium Financial Analysis & Strategic Insights"
    )
    
    # Create executive summary
    print("→ Generating executive summary...")
    summary_points = ai.generate_executive_summary(financial_data)
    prs.add_content_slide("Executive Summary", summary_points)
    
    # Create financial metrics slide
    print("→ Generating financial metrics...")
    prs.add_financial_data_slide(
        "Key Financial Metrics",
        financial_data,
        chart_type="metrics"
    )
    
    # Create detailed data table if available
    if 'table' in financial_data:
        print("→ Generating data table...")
        prs.add_financial_data_slide(
            "Financial Performance Details",
            financial_data,
            chart_type="table"
        )
    
    # Create comparison slide if available
    if 'comparisons' in financial_data:
        print("→ Generating comparison analysis...")
        prs.add_financial_data_slide(
            "Comparative Analysis",
            financial_data,
            chart_type="comparison"
        )
    
    # Create recommendations slide
    print("→ Generating strategic recommendations...")
    recommendations = ai.generate_recommendations(
        financial_data.get('analysis', {})
    )
    prs.add_content_slide("Strategic Recommendations", recommendations)
    
    # Save presentation
    print("\n→ Finalizing presentation...")
    output_path = prs.save(output_filename)
    
    print("=" * 60)
    print(f"✓ Quantum presentation complete!")
    print(f"  Output: {output_path}")
    print("=" * 60)
    
    return output_path


# ═══════════════════════════════════════════════════════════════════════════
# SECTION 2: USAGE EXAMPLES
# ═══════════════════════════════════════════════════════════════════════════

"""
Quantum Usage Examples
Demonstration of luxury financial presentation generation
"""


def example_1_quarterly_report():
    """Example: Quarterly Financial Report"""
    print("\n" + "="*60)
    print("EXAMPLE 1: Quarterly Financial Report")
    print("="*60)
    
    data = {
        'metrics': {
            'Total Revenue': '$245.8M',
            'Net Income': '$54.2M',
            'Profit Margin': '22.1%',
            'YoY Growth': '+18.5%',
            'EBITDA': '$72.3M',
            'Operating Cash Flow': '$68.9M'
        },
        'table': [
            ['Metric', 'Q4 2024', 'Q3 2024', 'Change'],
            ['Revenue', '$65.2M', '$61.8M', '+5.5%'],
            ['Gross Profit', '$42.1M', '$39.6M', '+6.3%'],
            ['Operating Income', '$18.7M', '$16.9M', '+10.7%'],
            ['Net Income', '$14.3M', '$13.1M', '+9.2%']
        ],
        'comparisons': {
            'Revenue per Employee': {'current': '$425K', 'change': '+12.3%'},
            'Customer Retention': {'current': '94.7%', 'change': '+2.1%'},
            'Market Penetration': {'current': '23.5%', 'change': '+4.2%'}
        },
        'analysis': {
            'trend': 'positive'
        }
    }
    
    return create_quantum_presentation(
        "Q4 2024 Quarterly Financial Report",
        data,
        output_filename="quantum_q4_report.pptx"
    )


def example_2_investment_pitch():
    """Example: Investment Pitch Deck"""
    print("\n" + "="*60)
    print("EXAMPLE 2: Investment Pitch Deck")
    print("="*60)
    
    data = {
        'metrics': {
            'Market Opportunity': '$12.5B',
            'Current Valuation': '$350M',
            'Revenue (TTM)': '$87.3M',
            'Growth Rate': '+156%',
            'Gross Margin': '68%',
            'Monthly Recurring Revenue': '$7.8M'
        },
        'comparisons': {
            'Customer Growth': {'current': '45,230', 'change': '+187%'},
            'Average Contract Value': {'current': '$24.5K', 'change': '+32%'},
            'Net Revenue Retention': {'current': '132%', 'change': '+8%'}
        },
        'analysis': {
            'trend': 'positive'
        }
    }
    
    return create_quantum_presentation(
        "Series B Investment Opportunity",
        data,
        output_filename="quantum_investment_pitch.pptx"
    )


def example_3_annual_review():
    """Example: Annual Performance Review"""
    print("\n" + "="*60)
    print("EXAMPLE 3: Annual Performance Review")
    print("="*60)
    
    data = {
        'metrics': {
            'Annual Revenue': '$892.4M',
            'Net Profit': '$178.9M',
            'ROE': '28.7%',
            'Dividend Yield': '3.2%',
            'EPS': '$4.85',
            'P/E Ratio': '24.3'
        },
        'table': [
            ['Year', '2024', '2023', '2022'],
            ['Revenue', '$892M', '$756M', '$628M'],
            ['Net Income', '$179M', '$143M', '$112M'],
            ['Total Assets', '$1.2B', '$1.0B', '$845M'],
            ['Shareholders Equity', '$623M', '$512M', '$428M']
        ],
        'comparisons': {
            'Operating Efficiency': {'current': '20.1%', 'change': '+3.5%'},
            'Asset Turnover': {'current': '1.43', 'change': '+0.12'},
            'Debt-to-Equity': {'current': '0.45', 'change': '-0.08'}
        },
        'analysis': {
            'trend': 'positive'
        }
    }
    
    return create_quantum_presentation(
        "2024 Annual Performance Review",
        data,
        output_filename="quantum_annual_review.pptx"
    )


def example_4_custom_styling():
    """Example: Custom Luxury Styling"""
    print("\n" + "="*60)
    print("EXAMPLE 4: Custom Luxury Styling")
    print("="*60)
    
    # Create custom luxury style
    custom_style = SlideStyle(
        title_font_size=48,
        title_font_name="Arial",
        title_color=(25, 25, 112),  # Midnight Blue
        body_font_size=20,
        body_font_name="Arial",
        body_color=(70, 70, 70),
        accent_color=(218, 165, 32)  # Goldenrod
    )
    
    # Create presentation with custom style
    prs = QuantumPresentation(custom_style)
    
    prs.add_title_slide(
        "Premium Portfolio Analysis",
        "Bespoke Financial Intelligence"
    )
    
    prs.add_content_slide(
        "Portfolio Highlights",
        [
            "Diversified asset allocation across 12 sectors",
            "Risk-adjusted returns exceeding benchmark by 4.2%",
            "Strategic rebalancing completed Q4 2024",
            "ESG compliance rating: AAA",
            "Liquidity position optimized for market conditions"
        ]
    )
    
    data = {
        'metrics': {
            'Portfolio Value': '$12.8M',
            'YTD Return': '+24.3%',
            'Sharpe Ratio': '1.85',
            'Alpha': '+4.2%',
            'Beta': '0.92',
            'Max Drawdown': '-8.1%'
        }
    }
    
    prs.add_financial_data_slide(
        "Performance Metrics",
        data,
        chart_type="metrics"
    )
    
    return prs.save("quantum_custom_style.pptx")


def example_5_with_reference_deck():
    """Example: Learning from Reference Deck"""
    print("\n" + "="*60)
    print("EXAMPLE 5: Learning from Reference Deck")
    print("="*60)
    print("Note: Place your reference deck at '/mnt/user-data/uploads/reference.pptx'")
    print("Quantum will analyze and learn the luxury formatting standards")
    print("="*60)
    
    reference_path = "/mnt/user-data/uploads/reference.pptx"
    
    data = {
        'metrics': {
            'Revenue': '$156.7M',
            'Growth': '+22.8%',
            'Margin': '31.2%',
            'EBITDA': '$48.9M'
        },
        'analysis': {
            'trend': 'positive'
        }
    }
    
    # This will automatically analyze the reference deck and apply its styling
    return create_quantum_presentation(
        "Financial Analysis with Learned Styles",
        data,
        reference_deck_path=reference_path,
        output_filename="quantum_with_reference.pptx"
    )


def run_all_examples():
    """Run all example demonstrations"""
    print("\n" + "🌟"*30)
    print("QUANTUM LUXURY FINANCIAL SOFTWARE")
    print("Complete Examples Demonstration")
    print("🌟"*30)
    
    examples = [
        ("Quarterly Report", example_1_quarterly_report),
        ("Investment Pitch", example_2_investment_pitch),
        ("Annual Review", example_3_annual_review),
        ("Custom Styling", example_4_custom_styling),
    ]
    
    results = []
    
    for name, example_func in examples:
        try:
            result = example_func()
            results.append((name, "✓ Success", result))
        except Exception as e:
            results.append((name, "✗ Error", str(e)))
    
    print("\n" + "="*60)
    print("SUMMARY OF GENERATED PRESENTATIONS")
    print("="*60)
    for name, status, result in results:
        print(f"{status} {name}: {result}")
    print("="*60)
    
    return results


# ═══════════════════════════════════════════════════════════════════════════
# SECTION 3: QUICK START SCRIPT
# ═══════════════════════════════════════════════════════════════════════════

def quantum_quickstart():
    """Generate a sample Quantum presentation - Simplest usage"""
    
    print("\n" + "="*70)
    print(" 🌟 QUANTUM - LUXURY FINANCIAL SOFTWARE TOOL")
    print("="*70)
    print("\nGenerating your first professional financial presentation...\n")
    
    # Sample financial data
    sample_data = {
        'metrics': {
            'Total Revenue': '$892.4M',
            'Net Income': '$178.9M',
            'Profit Margin': '20.1%',
            'YoY Growth': '+18.7%',
            'EBITDA': '$245.3M',
            'ROI': '28.4%'
        },
        'table': [
            ['Quarter', 'Revenue', 'Expenses', 'Net Income', 'Margin'],
            ['Q1 2024', '$215.2M', '$171.8M', '$43.4M', '20.2%'],
            ['Q2 2024', '$228.7M', '$182.4M', '$46.3M', '20.2%'],
            ['Q3 2024', '$231.9M', '$184.6M', '$47.3M', '20.4%'],
            ['Q4 2024', '$216.6M', '$173.1M', '$43.5M', '20.1%']
        ],
        'comparisons': {
            'Revenue Growth': {
                'current': '$892.4M',
                'change': '+18.7%'
            },
            'Operating Margin': {
                'current': '20.1%',
                'change': '+1.8%'
            },
            'Customer Base': {
                'current': '125,450',
                'change': '+32.5%'
            },
            'Market Share': {
                'current': '23.8%',
                'change': '+3.2%'
            }
        },
        'analysis': {
            'trend': 'positive'
        }
    }
    
    # Generate presentation
    output = create_quantum_presentation(
        title="2024 Annual Financial Performance",
        financial_data=sample_data,
        output_filename="quantum_quickstart_demo.pptx"
    )
    
    print("\n" + "="*70)
    print(f"✅ SUCCESS! Your presentation is ready: {output}")
    print("="*70)
    print("\n📋 What was generated:")
    print("   • Professional title slide")
    print("   • AI-generated executive summary")
    print("   • Key financial metrics dashboard")
    print("   • Detailed performance table")
    print("   • Comparative analysis with trends")
    print("   • Strategic recommendations")
    print("\n💡 Next Steps:")
    print("   1. Open the generated presentation")
    print("   2. Review the luxury formatting")
    print("   3. Upload your own reference deck to learn your brand")
    print("   4. Modify this script with your data")
    print("   5. Read documentation for advanced features")
    print("\n🎯 To use your own data, call create_quantum_presentation()")
    print("   with your financial data dictionary")
    print("\n" + "="*70 + "\n")
    
    return output


# ═══════════════════════════════════════════════════════════════════════════
# MAIN EXECUTION
# ═══════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    # Run the quick start by default
    quantum_quickstart()
    
    # Uncomment to run specific examples:
    # example_1_quarterly_report()
    # example_2_investment_pitch()
    # example_3_annual_review()
    # example_4_custom_styling()
    # run_all_examples()
